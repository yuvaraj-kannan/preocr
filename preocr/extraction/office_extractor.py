"""Office document extraction with structured output."""

from pathlib import Path
from typing import Dict, List, Optional, Any

from .. import exceptions
from ..utils.logger import get_logger
from .schemas import ExtractionResult, Element, Table, ElementType, BoundingBox, TableCell
from .base import generate_element_id, create_bbox, calculate_confidence

OfficeDocumentError = exceptions.OfficeDocumentError
logger = get_logger(__name__)

# Declare these as Optional[Any] so mypy knows they can be None
Document: Optional[Any]
Presentation: Optional[Any]
load_workbook: Optional[Any]

try:
    from docx import Document as _Document
    Document = _Document
except ImportError:
    Document = None

try:
    from pptx import Presentation as _Presentation
    Presentation = _Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook as _load_workbook
    load_workbook = _load_workbook
except ImportError:
    load_workbook = None


def extract_office_native_data(
    file_path: str,
    include_tables: bool = True,
    include_metadata: bool = True,
    include_structure: bool = True,
    include_bbox: bool = True,
) -> ExtractionResult:
    """
    Extract structured data from Office documents (DOCX, PPTX, XLSX).

    Args:
        file_path: Path to the office document
        include_tables: Whether to extract tables
        include_metadata: Whether to include document metadata
        include_structure: Whether to detect structure
        include_bbox: Whether to include bounding boxes (limited support)

    Returns:
        ExtractionResult with extracted data
    """
    path = Path(file_path)
    extension = path.suffix.lower()

    result = ExtractionResult(
        file_path=str(path),
        file_type=extension.lstrip("."),
        extraction_method="native",
        overall_confidence=0.0,
    )

    if extension == ".docx":
        return _extract_docx(path, result, include_tables, include_metadata, include_structure, include_bbox)
    elif extension == ".pptx":
        return _extract_pptx(path, result, include_metadata, include_structure, include_bbox)
    elif extension == ".xlsx":
        return _extract_xlsx(path, result, include_tables, include_metadata, include_structure, include_bbox)
    else:
        result.errors = [f"Unsupported office document type: {extension}"]
        return result


def _extract_docx(
    path: Path,
    result: ExtractionResult,
    include_tables: bool,
    include_metadata: bool,
    include_structure: bool,
    include_bbox: bool,
) -> ExtractionResult:
    """Extract from DOCX file."""
    if Document is None:
        result.errors = ["python-docx not available"]
        return result

    all_elements = []
    all_tables = []
    element_counter = 0
    table_counter = 0
    errors = []

    try:
        doc = Document(str(path))

        # Extract metadata
        if include_metadata:
            result.metadata.update({
                "extraction_method": "python-docx",
            })
            if doc.core_properties:
                props = doc.core_properties
                if props.title:
                    result.metadata["title"] = props.title
                if props.author:
                    result.metadata["author"] = props.author
                if props.created:
                    result.metadata["created"] = props.created.isoformat() if props.created else None

        # Extract paragraphs as elements
        for para_idx, paragraph in enumerate(doc.paragraphs):
            if not paragraph.text.strip():
                continue

            element_id = generate_element_id(f"elem_{element_counter}")
            element_counter += 1

            # Classify element type
            element_type = ElementType.NARRATIVE_TEXT
            if paragraph.style and paragraph.style.name:
                style_name = paragraph.style.name.lower()
                if "heading" in style_name or "title" in style_name:
                    element_type = ElementType.HEADING
                elif "title" in style_name:
                    element_type = ElementType.TITLE

            # Approximate bbox (limited support)
            if include_bbox:
                # Approximate position based on paragraph index
                y0 = para_idx * 20.0  # Approximate line height
                y1 = y0 + 15.0
                bbox = create_bbox(0, y0, 612, y1, 1, 612, 792)  # Standard letter size
            else:
                bbox = create_bbox(0, 0, 612, 792, 1, 612, 792)

            confidence = calculate_confidence(
                text_quality=0.9 if len(paragraph.text) > 10 else 0.7,
                extraction_method="python-docx",
                element_type_certainty=0.8,
                bbox_accuracy=0.5 if include_bbox else 0.3,
            )

            element = Element(
                element_id=element_id,
                element_type=element_type,
                text=paragraph.text,
                bbox=bbox,
                confidence=confidence,
                metadata={"style": paragraph.style.name if paragraph.style else None},
            )

            all_elements.append(element)

        # Extract tables
        if include_tables:
            for table_idx, table in enumerate(doc.tables):
                table_id = generate_element_id(f"table_{table_counter}")
                table_counter += 1

                cells = []
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if not cell_text:
                            continue

                        # Approximate cell bbox
                        if include_bbox:
                            cell_width = 612 / max(len(row.cells), 1)
                            cell_height = 20.0
                            cell_bbox = create_bbox(
                                col_idx * cell_width,
                                row_idx * cell_height,
                                (col_idx + 1) * cell_width,
                                (row_idx + 1) * cell_height,
                                1,
                                612,
                                792,
                            )
                        else:
                            cell_bbox = create_bbox(0, 0, 612, 792, 1, 612, 792)

                        cell_obj = TableCell(
                            row=row_idx,
                            col=col_idx,
                            text=cell_text,
                            bbox=cell_bbox,
                            confidence=0.9,
                        )

                        cells.append(cell_obj)

                if cells:
                    num_rows = len(table.rows)
                    num_cols = max(len(row.cells) for row in table.rows) if table.rows else 0

                    table_bbox = create_bbox(0, 0, 612, num_rows * 20.0, 1, 612, 792)

                    table_obj = Table(
                        element_id=table_id,
                        page_number=1,
                        bbox=table_bbox,
                        rows=num_rows,
                        columns=num_cols,
                        cells=cells,
                        confidence=0.9,
                        metadata={"extraction_method": "python-docx"},
                    )

                    all_tables.append(table_obj)

    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        errors.append(f"Extraction error: {str(e)}")

    # Populate result
    result.elements = all_elements
    result.tables = all_tables
    result.errors = errors

    # Calculate overall confidence
    all_confidences = [e.confidence for e in all_elements] + [t.confidence for t in all_tables]
    result.overall_confidence = sum(all_confidences) / len(all_confidences) if all_confidences else 0.0

    # Quality metrics
    result.quality_metrics = {
        "total_elements": len(all_elements),
        "total_tables": len(all_tables),
        "extraction_method": "python-docx",
    }

    return result


def _extract_pptx(
    path: Path,
    result: ExtractionResult,
    include_metadata: bool,
    include_structure: bool,
    include_bbox: bool,
) -> ExtractionResult:
    """Extract from PPTX file."""
    if Presentation is None:
        result.errors = ["python-pptx not available"]
        return result

    all_elements = []
    element_counter = 0
    errors = []

    try:
        prs = Presentation(str(path))

        # Extract metadata
        if include_metadata:
            result.metadata.update({
                "extraction_method": "python-pptx",
                "slide_count": len(prs.slides),
            })
            if prs.core_properties:
                props = prs.core_properties
                if props.title:
                    result.metadata["title"] = props.title
                if props.author:
                    result.metadata["author"] = props.author

        # Extract slides
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue

                element_id = generate_element_id(f"elem_{element_counter}")
                element_counter += 1

                # Approximate bbox
                if include_bbox and hasattr(shape, "left") and hasattr(shape, "top"):
                    bbox = create_bbox(
                        shape.left,
                        shape.top,
                        shape.left + (shape.width if hasattr(shape, "width") else 100),
                        shape.top + (shape.height if hasattr(shape, "height") else 100),
                        slide_idx + 1,
                        914400,  # PowerPoint units (1 inch = 914400 EMU)
                        685800,
                    )
                else:
                    bbox = create_bbox(0, 0, 914400, 685800, slide_idx + 1, 914400, 685800)

                element_type = ElementType.NARRATIVE_TEXT
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    element_type = ElementType.TITLE

                confidence = calculate_confidence(
                    text_quality=0.9 if len(shape.text) > 10 else 0.7,
                    extraction_method="python-pptx",
                    element_type_certainty=0.8,
                    bbox_accuracy=0.5 if include_bbox else 0.3,
                )

                element = Element(
                    element_id=element_id,
                    element_type=element_type,
                    text=shape.text,
                    bbox=bbox,
                    confidence=confidence,
                    metadata={"slide_number": slide_idx + 1},
                )

                all_elements.append(element)

    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        errors.append(f"Extraction error: {str(e)}")

    # Populate result
    result.elements = all_elements
    result.errors = errors

    # Calculate overall confidence
    all_confidences = [e.confidence for e in all_elements]
    result.overall_confidence = sum(all_confidences) / len(all_confidences) if all_confidences else 0.0

    # Quality metrics
    result.quality_metrics = {
        "total_elements": len(all_elements),
        "extraction_method": "python-pptx",
    }

    return result


def _extract_xlsx(
    path: Path,
    result: ExtractionResult,
    include_tables: bool,
    include_metadata: bool,
    include_structure: bool,
    include_bbox: bool,
) -> ExtractionResult:
    """Extract from XLSX file."""
    if load_workbook is None:
        result.errors = ["openpyxl not available"]
        return result

    all_tables = []
    table_counter = 0
    errors = []

    try:
        wb = load_workbook(str(path), data_only=True)

        # Extract metadata
        if include_metadata:
            result.metadata.update({
                "extraction_method": "openpyxl",
                "sheet_count": len(wb.sheetnames),
            })

        # Extract each sheet as a table
        if include_tables:
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                sheet = wb[sheet_name]

                table_id = generate_element_id(f"table_{table_counter}")
                table_counter += 1

                cells = []
                max_row = sheet.max_row
                max_col = sheet.max_column

                for row_idx in range(1, max_row + 1):
                    for col_idx in range(1, max_col + 1):
                        cell = sheet.cell(row_idx, col_idx)
                        cell_value = cell.value

                        if cell_value is None:
                            continue

                        cell_text = str(cell_value).strip()
                        if not cell_text:
                            continue

                        # Approximate cell bbox
                        if include_bbox:
                            cell_width = 100.0
                            cell_height = 20.0
                            cell_bbox = create_bbox(
                                (col_idx - 1) * cell_width,
                                (row_idx - 1) * cell_height,
                                col_idx * cell_width,
                                row_idx * cell_height,
                                sheet_idx + 1,
                                1000,
                                1000,
                            )
                        else:
                            cell_bbox = create_bbox(0, 0, 1000, 1000, sheet_idx + 1, 1000, 1000)

                        cell_obj = TableCell(
                            row=row_idx - 1,
                            col=col_idx - 1,
                            text=cell_text,
                            bbox=cell_bbox,
                            confidence=0.9,
                        )

                        cells.append(cell_obj)

                if cells:
                    table_bbox = create_bbox(0, 0, max_col * 100.0, max_row * 20.0, sheet_idx + 1, 1000, 1000)

                    table_obj = Table(
                        element_id=table_id,
                        page_number=sheet_idx + 1,
                        bbox=table_bbox,
                        rows=max_row,
                        columns=max_col,
                        cells=cells,
                        confidence=0.9,
                        metadata={
                            "extraction_method": "openpyxl",
                            "sheet_name": sheet_name,
                        },
                    )

                    all_tables.append(table_obj)

    except Exception as e:
        logger.error(f"XLSX extraction error: {e}")
        errors.append(f"Extraction error: {str(e)}")

    # Populate result
    result.tables = all_tables
    result.errors = errors

    # Calculate overall confidence
    all_confidences = [t.confidence for t in all_tables]
    result.overall_confidence = sum(all_confidences) / len(all_confidences) if all_confidences else 0.0

    # Quality metrics
    result.quality_metrics = {
        "total_tables": len(all_tables),
        "extraction_method": "openpyxl",
    }

    return result

